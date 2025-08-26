"""PPTX processing functionality for extracting notes and converting to images."""

import pandas as pd
import subprocess
import io
import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from typing import Dict, Any, Union
from pathlib import Path

from .file_utils import check_libreoffice, ensure_directory_exists
from ..config.settings import get_processing_config, get_image_quality_config


def extract_pptx_notes(pptx_path: Union[str, Path]) -> pd.DataFrame:
    """
    Extract speaker notes from PPTX file.

    Args:
        pptx_path: Path to PPTX file

    Returns:
        DataFrame with columns:
            - slide_number: Slide number (1-indexed)
            - slide_title: Title of the slide (first line of text, max 100 chars)
            - slide_text: Text content of the slide
            - speaker_notes: Speaker notes for the slide
            - has_notes: Whether the slide has speaker notes
            - notes_word_count: Number of words in the speaker notes
            - slide_text_word_count: Number of words in the slide text
            - image_filename: Suggested filename for the slide image (e.g. "slide-001.png")
    """
    pptx_path = Path(pptx_path)
    processing_config = get_processing_config()
    default_format = processing_config.get('default_format', 'png')

    prs = Presentation(pptx_path)
    slides_data = []

    for slide_num, slide in enumerate(prs.slides, 1):
        # Extract slide text (from shapes on the slide)
        slide_text = []
        slide_title = ""

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                slide_text.append(text)

                # Try to identify title (usually the first text or largest text)
                if not slide_title and text:
                    slide_title = text.split('\n')[0][:100]  # First line, max 100 chars

        # Extract speaker notes
        notes_text = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text.strip()

        slides_data.append({
            'slide_number': slide_num,
            'slide_title': slide_title,
            'slide_text': '\n'.join(slide_text),
            'speaker_notes': notes_text,
            'has_notes': bool(notes_text),
            'notes_word_count': len(notes_text.split()) if notes_text else 0,
            'slide_text_word_count': len('\n'.join(slide_text).split()) if slide_text else 0,
            'image_filename': f"slide-{slide_num:03d}.{default_format}"
        })

    return pd.DataFrame(slides_data)


def pptx_to_images_and_notes(
    pptx_path: Union[str, Path],
    output_dir: Union[str, Path],
    dpi: Union[int, None] = None,
    fmt: Union[str, None] = None,
    extract_notes: bool = True
) -> Dict[str, Any]:
    """
    Convert PPTX to images and extract notes.

    Args:
        pptx_path: Path to PPTX file
        output_dir: Directory to save images and notes
        dpi: Image resolution (uses config default if None)
        fmt: Image format - "png" or "jpeg" (uses config default if None)
        extract_notes: Whether to extract speaker notes

    Returns:
        Dict with image_files, notes_df, and output_dir

    Raises:
        FileNotFoundError: If PPTX file or LibreOffice not found
        RuntimeError: If conversion fails
    """
    pptx_path = Path(pptx_path).resolve()
    output_dir = Path(output_dir).resolve()

    # Load configuration
    processing_config = get_processing_config()
    image_quality_config = get_image_quality_config()

    if dpi is None:
        dpi = processing_config.get('default_dpi', 200)
    if fmt is None:
        fmt = processing_config.get('default_format', 'png')

    # Validate format
    supported_formats = processing_config.get('supported_formats', ['png', 'jpeg', 'jpg'])
    if fmt not in supported_formats:
        raise ValueError(f"Unsupported format '{fmt}'. Supported formats: {supported_formats}")

    ensure_directory_exists(output_dir)

    print(f"Processing: {pptx_path.name}")

    # Step 1: Extract notes if requested
    notes_df = pd.DataFrame()
    if extract_notes:
        print("Extracting speaker notes...")
        try:
            notes_df = extract_pptx_notes(pptx_path)
            # Update image filenames based on actual format
            notes_df['image_filename'] = notes_df['slide_number'].apply(
                lambda x: f"slide-{x:03d}.{fmt}"
            )

            notes_with_content = notes_df[notes_df['has_notes']]
            print(f"Found notes on {len(notes_with_content)} of {len(notes_df)} slides")

            # Save df to CSV
            csv_file = output_dir / f"{pptx_path.stem}_notes.csv"
            notes_df.to_csv(csv_file, index=False, encoding='utf-8')
            print(f"Notes df saved to: {csv_file}")

        except Exception as e:
            print(f"Warning: Could not extract notes: {e}")
            notes_df = pd.DataFrame()

    # Step 2: Convert to images
    soffice = check_libreoffice()

    # Clean previous outputs
    for f in output_dir.glob("slide*"):
        if f.is_file():
            f.unlink()

    # Convert PPTX → PDF
    pdf_path = output_dir / f"{pptx_path.stem}.pdf"

    print("Converting to PDF...")
    try:
        subprocess.run([
            soffice,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(output_dir),
            str(pptx_path)
        ], check=True, capture_output=True, text=True)
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"Failed to convert PPTX to PDF: {e}")

    if not pdf_path.exists():
        raise FileNotFoundError(f"PDF was not created: {pdf_path}")

    # Convert PDF → Images
    print(f"Converting to {fmt.upper()} images at {dpi} DPI...")
    doc = fitz.open(str(pdf_path))

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        zoom = dpi / 72
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)

        image_path = output_dir / f"slide-{page_num+1:03d}.{fmt}"

        if fmt.lower() == "png":
            pix.save(str(image_path))
        elif fmt.lower() in ["jpg", "jpeg"]:
            img_data = pix.tobytes("png")
            pil_img = Image.open(io.BytesIO(img_data))
            pil_img = pil_img.convert("RGB")

            # Get quality settings from config
            quality = image_quality_config.get('jpeg_quality', 90)
            optimize = image_quality_config.get('jpeg_optimize', True)

            pil_img.save(str(image_path), "JPEG", quality=quality, optimize=optimize)

        pix = None

    doc.close()

    # Summary
    image_files = sorted(list(output_dir.glob(f"slide*.{fmt}")))
    print(f"\nSuccessfully processed {len(image_files)} slides")
    print(f"Images saved to: {output_dir}")

    return {
        'image_files': image_files,
        'notes_df': notes_df,
        'output_dir': output_dir
    }
