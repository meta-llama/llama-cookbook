# Note: No logic for images yet
import os
from pptx import Presentation

class PPTParser:
    def __init__(self):
        pass
    
    def parse(self, file_path):
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        if not file_path.lower().endswith(('.pptx')):
            raise ValueError("Only .pptx format is supported")
        
        ppt = Presentation(file_path)
        all_text = []
        all_text.append(f"PowerPoint Presentation: {os.path.basename(file_path)}")
        all_text.append(f"Total Slides: {len(ppt.slides)}")
        all_text.append("")
        
        for i, slide in enumerate(ppt.slides):
            slide_num = i + 1
            all_text.append(f"Slide {slide_num}")
            all_text.append("-" * 40)
            
            if slide.shapes.title:
                all_text.append(f"Title: {slide.shapes.title.text}")
            
            slide_text = self._extract_slide_text(slide)
            if slide_text:
                all_text.append(slide_text)
                
            all_text.append("")
        
        return "\n".join(all_text)
    
    def _extract_slide_text(self, slide):
        texts = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = self._extract_text_frame(shape.text_frame)
                if text:
                    texts.append(text)
            
            # table logic
            if shape.has_table:
                table_text = self._extract_table(shape.table)
                if table_text:
                    texts.append(table_text)
        
        return "\n".join(texts)
    
    def _extract_text_frame(self, text_frame):
        text_lines = []
        
        for paragraph in text_frame.paragraphs:
            if paragraph.text.strip():
                text_lines.append(paragraph.text.strip())
        
        return "\n".join(text_lines)
    
    def _extract_table(self, table):
        table_lines = []
        
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                cell_text = ""
                for paragraph in cell.text_frame.paragraphs:
                    if paragraph.text.strip():
                        cell_text += paragraph.text.strip() + " "
                row_text.append(cell_text.strip())
            
            if any(row_text):  # Skip empty rows
                table_lines.append(" | ".join(row_text))
        
        return "\n".join(table_lines)
    
    def save(self, content, output_path):
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)